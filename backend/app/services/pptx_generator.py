from pathlib import Path

from pptx import Presentation

from app.services.parameter_calculator import image_slide_indices


def build_pptx(slides: list[dict], image_slides: int, output_path: Path) -> Path:
    prs = Presentation()
    image_indices = set(image_slide_indices(len(slides), image_slides))

    for idx, slide_data in enumerate(slides):
        layout = prs.slide_layouts[3] if idx in image_indices else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data.get("title", f"Slide {idx + 1}")

        if len(slide.placeholders) > 1:
            body = slide.placeholders[1].text_frame
            body.clear()
            for bullet in slide_data.get("bullets", []):
                p = body.add_paragraph()
                p.text = bullet

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
